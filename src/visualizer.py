import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np
import pandas as pd
from sklearn.metrics import confusion_matrix, ConfusionMatrixDisplay, roc_curve, auc, precision_recall_curve, average_precision_score
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import os
import shap 
from PIL import Image 

# 从 src 包中导入 config 模块
from src import config

def set_chinese_font():
    """
    设置 Matplotlib 字体以支持中文显示。
    优先使用 'SimHei',如果不可用则尝试 'Microsoft YaHei'。
    """
    try:
        plt.rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei', 'Arial Unicode MS']
        plt.rcParams['axes.unicode_minus'] = False # 解决负号显示问题
        print("Matplotlib 字体设置为 SimHei 或 Microsoft YaHei 以支持中文。")
    except Exception as e:
        print(f"设置中文 Matplotlib 字体失败: {e}。请确保已安装中文字体。")
        # 回退到默认字体,但可能无法显示中文
        plt.rcParams['font.sans-serif'] = ['DejaVu Sans']
        plt.rcParams['axes.unicode_minus'] = True

def save_figure(fig, filename: str, experiment_name: str, subdirectory: str = "") -> Path:
    """
    保存 Matplotlib 图形到指定路径。
    文件路径结构: reports_output_dir / experiment_name / subdirectory / filename

    Args:
        fig (matplotlib.figure.Figure): 要保存的 Matplotlib Figure 对象。
        filename (str): 保存的文件名 (例如 "my_plot.png")。
        experiment_name (str): 当前实验的名称,用于创建实验子目录。
        subdirectory (str): 实验子目录下的进一步子目录 (例如 "supervised_models", "shap_plots")。

    Returns:
        Path: 保存的文件的完整路径。
    """
    save_dir = config.REPORTS_OUTPUT_DIR / experiment_name
    if subdirectory:
        save_dir = save_dir / subdirectory
    
    save_dir.mkdir(parents=True, exist_ok=True) # 确保目录存在
    file_path = save_dir / filename
    
    fig.savefig(file_path, bbox_inches='tight', dpi=300)
    plt.close(fig) # 关闭图形,释放内存
    print(f"图表已保存: {file_path}")
    return file_path


def plot_and_save_figure(plot_func, filename: str, experiment_name: str, subdirectory: str = "", **kwargs) -> Path:
    """
    通用函数,用于调用绘图函数生成图形,然后保存并关闭图形。

    Args:
        plot_func (function): 一个接受特定参数并返回 matplotlib.figure.Figure 对象的绘图函数。
                                例如: plot_confusion_matrix, plot_feature_importance, plot_shap_summary 等。
        filename (str): 保存的文件名 (例如 "confusion_matrix.png")。
        experiment_name (str): 当前实验的名称,用于创建保存路径的根目录。
        subdirectory (str): 在 experiment_name 目录下的子目录,用于进一步组织文件。
        **kwargs: 传递给 plot_func 的任意关键字参数。

    Returns:
        Path: 保存的图表的完整文件路径,如果绘图失败则返回 None。
    """
    fig = None
    file_path = None
    try:
        fig = plot_func(**kwargs)
        if fig: # 确保图表对象不是 None
            file_path = save_figure(fig, filename, experiment_name, subdirectory)
        else:
            print(f"警告: 绘图函数 {plot_func.__name__} 未成功生成图形,未保存文件 {filename}。")
    except Exception as e:
        print(f"绘制或保存图表 '{filename}' 失败 (函数: {plot_func.__name__}): {e}")
        import traceback
        traceback.print_exc()
    finally:
        # save_figure 内部已经关闭了 fig，这里不再需要额外的 plt.close(fig)
        # 除非 fig 是在 save_figure 外部创建但未传入的情况，但这不符合当前设计
        pass 
    return file_path


def plot_confusion_matrix(y_true, y_pred, labels, title="混淆矩阵", cmap=plt.cm.Blues):
    """
    绘制混淆矩阵。

    Args:
        y_true (array-like): 真实标签。
        y_pred (array-like): 预测标签。
        labels (list): 标签类别列表 (例如 [0, 1])。
        title (str): 图表的标题。
        cmap (matplotlib.colors.Colormap): 颜色映射。

    Returns:
        matplotlib.figure.Figure: 混淆矩阵图的 Figure 对象。
    """
    cm = confusion_matrix(y_true, y_pred, labels=labels)
    disp = ConfusionMatrixDisplay(confusion_matrix=cm, display_labels=labels)
    fig, ax = plt.subplots(figsize=(8, 6))
    disp.plot(cmap=cmap, ax=ax, values_format='d')
    ax.set_title(title)
    return fig

def plot_feature_importance(feature_importances: np.ndarray, feature_names: list, top_n: int = 20, title: str = "特征重要性"):
    """
    绘制特征重要性柱状图。

    Args:
        feature_importances (np.ndarray): 特征重要性数值数组。
        feature_names (list): 特征名称列表。
        top_n (int): 显示前 N 个最重要的特征。
        title (str): 图表的标题。

    Returns:
        matplotlib.figure.Figure: 特征重要性图的 Figure 对象。
    """
    if not isinstance(feature_importances, np.ndarray) or feature_importances.size == 0 or len(feature_names) == 0:
        print("警告: 无特征重要性数据或特征名称。")
        return None
    
    if len(feature_importances) != len(feature_names):
        print("错误: 特征重要性数组和特征名称列表长度不匹配。")
        return None

    feature_importances_series = pd.Series(feature_importances, index=feature_names).sort_values(ascending=False)
    
    if top_n > 0:
        feature_importances_series = feature_importances_series.head(top_n)

    fig, ax = plt.subplots(figsize=(10, max(6, len(feature_importances_series) * 0.4)))
    sns.barplot(x=feature_importances_series.values, y=feature_importances_series.index, ax=ax, color="skyblue")
    ax.set_title(title)
    ax.set_xlabel("重要性得分")
    ax.set_ylabel("特征")
    plt.tight_layout()
    return fig

def plot_shap_summary(explainer=None, X=None, feature_names=None, plot_type: str = "bar", 
                      title: str = "SHAP Summary Plot", shap_values=None, expected_value=None, top_n: int = 20, sort: bool = True):
    """
    生成 SHAP 汇总图 (条形图或散点图)。
    可以传入 explainer 和 X,或者传入预计算的 shap_values 和 expected_value。

    Args:
        explainer: 已训练的 SHAP explainer 对象 (shap.Explainer, shap.TreeExplainer 等)。
        X (pd.DataFrame or np.ndarray): 用于生成 SHAP 值的特征数据。
        feature_names (list): 特征名称列表。
        plot_type (str): "bar" (平均绝对 SHAP 值) 或 "dot" (beeswarm 散点图)。
        title (str): 图表标题。
        shap_values (np.array or list): 预计算的 SHAP 值。
        expected_value (float or list): 预计算的 SHAP 期望值。
        top_n (int): 显示前N个特征，None为全部。
        sort (bool): 是否按重要性排序。
    
    Returns:
        matplotlib.figure.Figure: SHAP 汇总图的 Figure 对象。
    """
    if explainer is None and (shap_values is None or expected_value is None):
        print("错误: 必须提供 explainer 和 X,或预计算的 shap_values 和 expected_value。")
        return None

    try:
        shap_values_to_plot = None
        X_to_plot = X if isinstance(X, pd.DataFrame) else pd.DataFrame(X, columns=feature_names)

        if explainer is not None and X is not None:
            if hasattr(explainer, 'shap_values'):
                shap_values_to_plot = explainer.shap_values(X_to_plot)
                expected_value_to_plot = explainer.expected_value
            elif isinstance(explainer, shap.Explainer):
                e_values = explainer(X_to_plot)
                shap_values_to_plot = e_values.values
                expected_value_to_plot = e_values.base_values
            else:
                print("警告: 提供的 explainer 类型不支持直接计算 SHAP 值,请检查。")
                return None
            if isinstance(shap_values_to_plot, list) and len(shap_values_to_plot) == 2:
                shap_values_to_plot = shap_values_to_plot[1]
                if isinstance(expected_value_to_plot, list) and len(expected_value_to_plot) == 2:
                    expected_value_to_plot = expected_value_to_plot[1]
        elif shap_values is not None and expected_value is not None:
            shap_values_to_plot = shap_values
            expected_value_to_plot = expected_value
        else:
            print("错误: 无法获取用于绘图的 SHAP 值。")
            return None

        if shap_values_to_plot is None or np.all(shap_values_to_plot == 0):
            print("警告: SHAP 值全为零或为空，无法绘制 SHAP 汇总图。这可能表明模型没有学习到任何有用的特征模式。")
            return None

        if feature_names is None:
            if isinstance(X_to_plot, pd.DataFrame):
                feature_names = X_to_plot.columns.tolist()
            else:
                feature_names = [f"Feature {i}" for i in range(X_to_plot.shape[1])]

        plt.clf() 
        plt.figure(figsize=(10, 6))

        if plot_type == "bar":
            # 计算平均绝对SHAP值
            mean_abs_shap = np.abs(shap_values_to_plot).mean(axis=0)
            shap_series = pd.Series(mean_abs_shap, index=feature_names)
            if sort:
                shap_series = shap_series.sort_values(ascending=False)
            if top_n is not None and top_n > 0:
                shap_series = shap_series.head(top_n)
            # 画bar图
            fig, ax = plt.subplots(figsize=(10, max(6, len(shap_series) * 0.4)))
            sns.barplot(x=shap_series.values, y=shap_series.index, ax=ax, color="skyblue")
            ax.set_title(title)
            ax.set_xlabel("平均绝对SHAP值")
            ax.set_ylabel("特征")
            plt.tight_layout()
            return fig
        elif plot_type == "dot":
            shap.summary_plot(shap_values_to_plot, X_to_plot, feature_names=feature_names, 
                              plot_type="dot", show=False, max_display=config.SHAP_MAX_DISPLAY)
            plt.title(title)
            plt.tight_layout()
            return plt.gcf()
        else:
            print(f"不支持的 SHAP 绘制类型: {plot_type}")
            plt.close(plt.gcf())
            return None
    except Exception as e:
        print(f"plot_shap_summary() 内部错误: {e}")
        import traceback
        traceback.print_exc()
        return None
    finally:
        pass


def plot_shap_dependence(explainer, X, feature_name: str, feature_names: list, interaction_index=None, title: str = "SHAP 依赖图"):
    """
    绘制 SHAP 依赖图。

    Args:
        explainer: SHAP Explainer 对象。
        X (pd.DataFrame or np.ndarray): 用于解释的数据 (通常是测试集)。
        feature_name (str): 要绘制依赖图的特征名称。
        feature_names (list): 所有特征的名称列表。
        interaction_index (str or None): 用于着色的交互特征名称。
        title (str): 图表标题。

    Returns:
        matplotlib.figure.Figure: SHAP 依赖图的 Figure 对象。
    """
    try:
        # 确保X的列名是特征名,这是shap.dependence_plot的要求
        if not isinstance(X, pd.DataFrame):
            X_df = pd.DataFrame(X, columns=feature_names) 
        else:
            X_df = X

        shap_values = explainer.shap_values(X_df)
        if isinstance(shap_values, list) and len(shap_values) == 2:
            shap_values = shap_values[1] # 取正类 (1) 的 SHAP 值

        fig, ax = plt.subplots(figsize=(10, 7))
        shap.dependence_plot(
            feature_name, 
            shap_values, 
            X_df, 
            interaction_index=interaction_index, 
            show=False, 
            ax=ax
        )
        ax.set_title(title)
        plt.tight_layout()
        return fig
    except Exception as e:
        print(f"绘制 SHAP 依赖图失败 for {feature_name}: {e}")
        import traceback
        traceback.print_exc()
        return None

def plot_cluster_scatter(reduced_features_df: pd.DataFrame, x_col: str, y_col: str, cluster_col: str, title: str, legend_title: str = "Cluster"):
    """
    绘制聚类结果的二维散点图。

    Args:
        reduced_features_df (pd.DataFrame): 包含降维特征和聚类标签的DataFrame。
        x_col (str): X轴的列名 (例如 "Component_1")。
        y_col (str): Y轴的列名 (例如 "Component_2")。
        cluster_col (str): 聚类标签的列名 (例如 "Cluster")。
        title (str): 图表标题。
        legend_title (str): 图例标题。

    Returns:
        matplotlib.figure.Figure: 聚类散点图的 Figure 对象。
    """
    fig, ax = plt.subplots(figsize=(10, 8))
    sns.scatterplot(
        x=x_col, 
        y=y_col, 
        hue=cluster_col, 
        data=reduced_features_df, 
        color="skyblue", 
        s=100, 
        ax=ax
    )
    ax.set_title(title)
    ax.set_xlabel(x_col)
    ax.set_ylabel(y_col)
    ax.legend(title=legend_title, bbox_to_anchor=(1.05, 1), loc='upper left')
    plt.tight_layout()
    return fig

def plot_elbow_method(wcss_values: list, max_k: int, title: str = "KMeans 肘部法则"):
    """
    绘制 KMeans 肘部法则图。

    Args:
        wcss_values (list): 不同 k 值对应的 WCSS (簇内平方和) 值列表。
        max_k (int): 计算 WCSS 的最大簇数量 k。
        title (str): 图表标题。

    Returns:
        matplotlib.figure.Figure: 肘部法则图的 Figure 对象。
    """
    fig, ax = plt.subplots(figsize=(8, 6))
    sns.lineplot(x=range(1, max_k + 1), y=wcss_values, marker='o', ax=ax)
    ax.set_title(title)
    ax.set_xlabel("聚类数量 (K)")
    ax.set_ylabel("簇内平方和 (WCSS)")
    ax.set_xticks(range(1, max_k + 1))
    plt.grid(True)
    plt.tight_layout()
    return fig

def generate_powerpoint_report(figure_paths_with_titles: list, ppt_filename: str, experiment_name: str):
    """
    生成包含所有图表的 PowerPoint 报告。

    Args:
        figure_paths_with_titles (list): 包含 (文件路径, 图表标题) 元组的列表。
                                          文件路径应为 Path 对象。
        ppt_filename (str): PowerPoint 文件的名称 (例如 "report.pptx")。
        experiment_name (str): 当前实验的名称,用于报告目录。
    """
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6] # 空白布局

    # 确保报告保存目录存在
    report_dir = config.REPORTS_OUTPUT_DIR / experiment_name
    report_dir.mkdir(parents=True, exist_ok=True)
    output_ppt_path = report_dir / ppt_filename

    for img_path, title_text in figure_paths_with_titles:
        if img_path.exists():
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加标题文本框
            left = Inches(1)
            top = Inches(0.5)
            width = Inches(8)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = title_text
            p.font.size = Inches(0.3) # 适当调整字体大小
            
            # 添加图片并保持宽高比
            try:
                img = Image.open(img_path)
                img_width_px, img_height_px = img.size
                aspect_ratio = img_width_px / img_height_px

                # 定义幻灯片区域（在标题下方）
                slide_width_inches = prs.slide_width.inches
                slide_height_inches = prs.slide_height.inches
                
                # 图片可用的最大宽度和高度 (留出边距)
                max_img_width = Inches(slide_width_inches - 2) # 左右各1英寸边距
                max_img_height = Inches(slide_height_inches - top.inches - 1) # 标题下方，底部1英寸边距

                # 计算图片在幻灯片上的实际尺寸
                if img_width_px > img_height_px: # 宽图
                    pic_width = min(max_img_width, Inches(img_width_px * (max_img_height.inches / img_height_px)))
                    pic_height = pic_width / aspect_ratio
                else: # 高图或方图
                    pic_height = min(max_img_height, Inches(img_height_px * (max_img_width.inches / img_width_px)))
                    pic_width = pic_height * aspect_ratio

                # 居中放置图片
                left = (prs.slide_width - pic_width) / 2
                top = (prs.slide_height - pic_height) / 2 + Inches(0.5) # 标题下方，稍微向下偏移
                
                pic = slide.shapes.add_picture(str(img_path), left, top, width=pic_width, height=pic_height)
            except Exception as e:
                print(f"无法将图片 {img_path} 添加到PPT: {e}")
        else:
            print(f"警告: 图片文件不存在,无法添加到报告: {img_path}")

    prs.save(output_ppt_path)
    print(f"PowerPoint 报告已生成: {output_ppt_path}")


def display_performance_metrics(model_name: str, metrics: dict):
    """
    将模型的性能指标以美观的表格格式打印出来,并自动调整列宽。

    Args:
        model_name (str): 模型的名称。
        metrics (dict): 包含性能指标的字典。
    """
    print(f"\n--- {model_name} 模型性能摘要 ---")

    # 定义希望显示的指标顺序和友好名称
    metric_order = {
        "accuracy": "准确率",
        "precision_weighted": "加权精确率",
        "recall_weighted": "加权召回率",
        "f1_weighted": "加权F1分数",
        "roc_auc": "ROC AUC",
        "Class_0_precision": "类别 0 精确率",
        "Class_0_recall": "类别 0 召回率",
        "Class_0_f1-score": "类别 0 F1分数",
        "Class_0_support": "类别 0 数据量",
        "Class_1_precision": "类别 1 精确率",
        "Class_1_recall": "类别 1 召回率",
        "Class_1_f1-score": "类别 1 F1分数",
        "Class_1_support": "类别 1 数据量",
        # 可以根据需要添加无监督指标
        "silhouette_score": "轮廓系数",
        "calinski_harabasz_score": "Calinski-Harabasz 分数"
    }

    # 准备用于 DataFrame 的数据
    data_rows = []
    
    # 提取并组织数据
    for key, display_name in metric_order.items():
        if key in metrics: # 检查键是否存在于 metrics 字典中
            value = metrics[key]
            # 格式化数值,保留小数点后4位,支持整数
            # 检查值是否为浮点数，如果是，保留4位小数；否则，如果为整数，则直接转换为整数，否则转为字符串
            formatted_value = f"{value:.4f}" if isinstance(value, (float, np.floating)) else (
                               f"{int(value)}" if isinstance(value, (int, np.integer)) else str(value))
            data_rows.append([display_name, formatted_value])
            
    # 创建 DataFrame
    display_df = pd.DataFrame(data_rows, columns=["指标", "值"])

    if display_df.empty:
        print("没有可显示的性能指标。")
        return

    # 计算列宽
    # 列标题长度
    col_header_len_metric = len(display_df.columns[0])
    col_header_len_value = len(display_df.columns[1])

    # 数据内容的最大长度
    max_len_metric = display_df["指标"].str.len().max()
    max_len_value = display_df["值"].str.len().max()

    # 最终列宽：取标题和内容中的最大值,再加2作为填充
    width_metric = max(col_header_len_metric, max_len_metric) + 2
    width_value = max(col_header_len_value, max_len_value) + 2

    # 打印头部
    header_metric = f"{display_df.columns[0]:<{width_metric}}"
    header_value = f"{display_df.columns[1]:<{width_value}}"
    print(f"{header_metric}{header_value}")
    print("-" * (width_metric + width_value))

    # 打印数据行
    for _, row in display_df.iterrows():
        metric_str = f"{row['指标']:<{width_metric}}"
        value_str = f"{row['值']:<{width_value}}"
        print(f"{metric_str}{value_str}")
    
    print("-" * (width_metric + width_value))


def plot_roc_pr_curve(y_true, y_pred_proba, plot_type: str = "ROC", title: str = ""):
    """
    绘制 ROC 或 Precision-Recall 曲线。

    Args:
        y_true (array-like): 真实标签。
        y_pred_proba (array-like): 预测为正类的概率。
        plot_type (str): "ROC" 或 "PR"。
        title (str): 图表标题。

    Returns:
        matplotlib.figure.Figure: 曲线图的 Figure 对象。
    """
    fig, ax = plt.subplots(figsize=(8, 8))

    if plot_type == "ROC":
        fpr, tpr, _ = roc_curve(y_true, y_pred_proba)
        roc_auc = auc(fpr, tpr)
        ax.plot(fpr, tpr, color='darkorange', lw=2, label=f'ROC curve (area = {roc_auc:.2f})')
        ax.plot([0, 1], [0, 1], color='navy', lw=2, linestyle='--')
        ax.set_xlim([0.0, 1.0])
        ax.set_ylim([0.0, 1.05])
        ax.set_xlabel('False Positive Rate')
        ax.set_ylabel('True Positive Rate')
        default_title = "ROC 曲线"
    elif plot_type == "PR":
        precision, recall, _ = precision_recall_curve(y_true, y_pred_proba)
        avg_precision = average_precision_score(y_true, y_pred_proba)
        ax.plot(recall, precision, color='navy', lw=2, label=f'PR curve (area = {avg_precision:.2f})')
        ax.set_xlim([0.0, 1.05])
        ax.set_ylim([0.0, 1.05])
        ax.set_xlabel('Recall')
        ax.set_ylabel('Precision')
        default_title = "Precision-Recall 曲线"
    else:
        raise ValueError("plot_type 必须是 'ROC' 或 'PR'")

    ax.set_title(title if title else default_title)
    ax.legend(loc="lower right")
    plt.tight_layout()
    return fig